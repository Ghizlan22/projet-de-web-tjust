from fastapi import FastAPI, File, UploadFile, Form
from fastapi.responses import JSONResponse
from transformers import pipeline
from typing import Optional
import io
from PIL import Image
import tempfile
import os
import fitz  # PyMuPDF
import docx
import pandas as pd
import pptx
from fastapi.middleware.cors import CORSMiddleware
from langdetect import detect

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Liste des langues supportées
SUPPORTED_LANGUAGES = ["fr", "en", "de", "es", "it", "zh", "ar"]

# Modèles de traduction valides (existants sur Hugging Face)
translation_models = {
    "fr-en": "Helsinki-NLP/opus-mt-fr-en",
    "en-fr": "Helsinki-NLP/opus-mt-en-fr",
    "fr-de": "Helsinki-NLP/opus-mt-fr-de",
    "de-fr": "Helsinki-NLP/opus-mt-de-fr",
    "fr-es": "Helsinki-NLP/opus-mt-fr-es",
    "es-fr": "Helsinki-NLP/opus-mt-es-fr",
    "en-zh": "Helsinki-NLP/opus-mt-en-zh",
    "zh-en": "Helsinki-NLP/opus-mt-zh-en",
    "en-it": "Helsinki-NLP/opus-mt-en-it",
    "it-en": "Helsinki-NLP/opus-mt-it-en",
    "en-ar": "Helsinki-NLP/opus-mt-en-ar",
    "ar-en": "Helsinki-NLP/opus-mt-ar-en",
    "en-es": "Helsinki-NLP/opus-mt-en-es",
    "en-de": "Helsinki-NLP/opus-mt-en-de",
    "es-ar": "Helsinki-NLP/opus-mt-es-ar",
    "es-en": "Helsinki-NLP/opus-mt-es-en",
    "es-de": "Helsinki-NLP/opus-mt-es-de",
    "es-it": "Helsinki-NLP/opus-mt-es-it",
    "es-zh": "Helsinki-NLP/opus-mt-es-zh",
    "ar-fr": "Helsinki-NLP/opus-mt-ar-fr",
    "ar-de": "Helsinki-NLP/opus-mt-ar-de",
    "ar-es": "Helsinki-NLP/opus-mt-ar-es",
    "ar-it": "Helsinki-NLP/opus-mt-ar-it",
    "ar-zh": "Helsinki-NLP/opus-mt-ar-zh",
    "de-en": "Helsinki-NLP/opus-mt-de-en",
    "de-de": "Helsinki-NLP/opus-mt-de-de",
    "de-es": "Helsinki-NLP/opus-mt-de-es",
    "de-it": "Helsinki-NLP/opus-mt-de-it",
    "de-zh": "Helsinki-NLP/opus-mt-de-zh",
    "de-ar": "Helsinki-NLP/opus-mt-de-ar",
    "it-fr": "Helsinki-NLP/opus-mt-it-fr",
    "it-de": "Helsinki-NLP/opus-mt-it-de",
    "it-es": "Helsinki-NLP/opus-mt-it-es",
    "it-zh": "Helsinki-NLP/opus-mt-it-zh",
    "it-ar": "Helsinki-NLP/opus-mt-it-ar",
    "zh-fr": "Helsinki-NLP/opus-mt-zh-fr",
    "zh-de": "Helsinki-NLP/opus-mt-zh-en",
    "zh-it": "Helsinki-NLP/opus-mt-zh-it",
    "zh-es": "Helsinki-NLP/opus-mt-zh-es",
    "zh-ar": "Helsinki-NLP/opus-mt-zh-ar",
    

}

def extract_text_from_pdf(file_path):
    text = ""
    with fitz.open(file_path) as doc:
        for page in doc:
            text += page.get_text("text") + "\n"
    return text

def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pptx(file_path):
    presentation = pptx.Presentation(file_path)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_excel(file_path):
    df = pd.read_excel(file_path, engine="openpyxl")
    return df.to_string(index=False)

def chunk_text(text, max_length=512):
    words = text.split()
    chunks, current_chunk = [], []

    for word in words:
        if len(" ".join(current_chunk) + " " + word) <= max_length:
            current_chunk.append(word)
        else:
            chunks.append(" ".join(current_chunk))
            current_chunk = [word]

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    return chunks

def translate_text(text, source_lang, target_lang):
    if source_lang not in SUPPORTED_LANGUAGES or target_lang not in SUPPORTED_LANGUAGES:
        return None  # Langue non supportée

    model_key = f"{source_lang}-{target_lang}"
    if model_key in translation_models:
        model_name = translation_models[model_key]
        translator = pipeline("translation", model=model_name)
        translated_chunks = [translator(chunk)[0]["translation_text"] for chunk in chunk_text(text)]
        return " ".join(translated_chunks)

    # Si pas de traduction directe, utiliser l'anglais comme pivot
    model_to_en = f"{source_lang}-en"
    model_from_en = f"en-{target_lang}"

    if model_to_en in translation_models and model_from_en in translation_models:
        translator_to_en = pipeline("translation", model=translation_models[model_to_en])
        translator_from_en = pipeline("translation", model=translation_models[model_from_en])

        intermediate_texts = [translator_to_en(chunk)[0]["translation_text"] for chunk in chunk_text(text)]
        intermediate_text = " ".join(intermediate_texts)

        final_texts = [translator_from_en(chunk)[0]["translation_text"] for chunk in chunk_text(intermediate_text)]
        return " ".join(final_texts)

    return None  # Pas de modèle disponible

@app.post("/translate")
async def translate_document(file: UploadFile = File(...), language: str = Form(...)):
    try:
        suffix = file.filename.split(".")[-1].lower()
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=f".{suffix}")
        temp_file.write(await file.read())
        temp_file.close()

        extractors = {
            "pdf": extract_text_from_pdf,
            "docx": extract_text_from_docx,
            "pptx": extract_text_from_pptx,
            "xls": extract_text_from_excel,
            "xlsx": extract_text_from_excel
        }

        if suffix not in extractors:
            return JSONResponse({"error": "Format non supporté"}, status_code=400)

        text = extractors[suffix](temp_file.name)
        os.remove(temp_file.name)

        if not text.strip():
            return JSONResponse({"error": "Aucun texte détecté"}, status_code=400)

        detected_lang = detect(text)
        if detected_lang not in SUPPORTED_LANGUAGES:
            return JSONResponse({"error": f"Langue non supportée : {detected_lang}"}, status_code=400)

        if detected_lang == language:
            return JSONResponse({"translation": text, "note": "Déjà dans la langue choisie."})

        translated_text = translate_text(text, detected_lang, language)
        if translated_text:
            return JSONResponse({"translation": translated_text})
        else:
            return JSONResponse({"error": "Aucun modèle de traduction trouvé."}, status_code=400)

    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)
